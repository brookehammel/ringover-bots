import streamlit as st
import gspread
from google.oauth2.service_account import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import anthropic
import pandas as pd
import json
import csv
import io
import base64
import os
import re

# ============================================================
# RINGOVER BOTS
# Book of Business Bot & Process Bot
# ============================================================

# --- Page Configuration ---
st.set_page_config(
    page_title="Ringover Bots",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- Ringover Brand Colors ---
NAVY = "#001b39"
TEAL = "#27c9d6"
WHITE = "#ffffff"

# --- Styling (Ringover branding with Poppins font) ---
st.markdown(f"""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700;800&display=swap');

    html, body, [class*="css"], .stApp, .stMarkdown, .stTextInput, .stButton {{
        font-family: 'Poppins', sans-serif !important;
    }}

    .stApp {{
        background-color: {NAVY};
        color: {WHITE};
    }}

    /* Sidebar */
    [data-testid="stSidebar"] {{
        background-color: #00122a;
    }}
    [data-testid="stSidebar"] * {{
        color: {WHITE} !important;
    }}

    /* Text input */
    .stTextInput input {{
        background-color: #00122a;
        color: {WHITE};
        border: 1px solid {TEAL};
        border-radius: 6px;
    }}

    /* Chat input — light background */
    [data-testid="stChatInput"] {{
        background-color: {WHITE} !important;
        border-radius: 8px;
    }}
    [data-testid="stChatInput"] > div {{
        background-color: {WHITE} !important;
        border: 1px solid {TEAL} !important;
        border-radius: 8px;
    }}
    .stChatInput textarea, [data-testid="stChatInput"] textarea {{
        background-color: {WHITE} !important;
        color: {NAVY} !important;
        border: none !important;
        caret-color: {NAVY} !important;
    }}
    .stChatInput textarea::placeholder, [data-testid="stChatInput"] textarea::placeholder {{
        color: #8a9ba8 !important;
    }}
    /* Send button inside chat input */
    [data-testid="stChatInput"] button {{
        color: {NAVY} !important;
    }}
    [data-testid="stChatInput"] button svg {{
        fill: {NAVY} !important;
    }}

    /* Default Streamlit buttons */
    .stButton > button {{
        background-color: {TEAL};
        color: {NAVY};
        font-family: 'Poppins', sans-serif;
        font-weight: 600;
        border: none;
        border-radius: 8px;
        padding: 0.6rem 1.4rem;
        transition: all 0.2s ease;
    }}
    .stButton > button:hover {{
        background-color: {WHITE};
        color: {NAVY};
        transform: translateY(-1px);
    }}

    /* Chat messages */
    [data-testid="stChatMessage"] {{
        background-color: #002a4f !important;
        border-radius: 10px;
        padding: 0.5rem 1rem;
    }}
    [data-testid="stChatMessage"] *,
    [data-testid="stChatMessage"] p,
    [data-testid="stChatMessage"] li,
    [data-testid="stChatMessage"] span,
    [data-testid="stChatMessage"] strong,
    [data-testid="stChatMessage"] em,
    [data-testid="stChatMessage"] h1,
    [data-testid="stChatMessage"] h2,
    [data-testid="stChatMessage"] h3,
    [data-testid="stChatMessage"] h4,
    [data-testid="stChatMessage"] h5,
    [data-testid="stChatMessage"] h6,
    [data-testid="stChatMessage"] td,
    [data-testid="stChatMessage"] th {{
        color: {WHITE} !important;
    }}
    [data-testid="stChatMessage"] a {{
        color: {TEAL} !important;
    }}
    [data-testid="stChatMessage"] code {{
        background-color: #00122a !important;
        color: {TEAL} !important;
        padding: 2px 6px;
        border-radius: 4px;
    }}
    [data-testid="stChatMessage"] pre code {{
        background-color: #00122a !important;
        color: {WHITE} !important;
    }}

    /* Main content area text — force white on navy (aggressive overrides) */
    html body .main .block-container,
    html body .main .block-container div,
    html body .main .block-container p,
    html body .main .block-container li,
    html body .main .block-container span,
    html body .main .block-container strong,
    html body .main .block-container em,
    html body .main .block-container h1,
    html body .main .block-container h2,
    html body .main .block-container h3,
    html body .main .block-container h4,
    html body .main .block-container h5,
    html body .main .block-container h6,
    html body .main .block-container td,
    html body .main .block-container th,
    html body .main .block-container label,
    html body .stMarkdown,
    html body .stMarkdown *,
    html body [data-testid="stMarkdownContainer"],
    html body [data-testid="stMarkdownContainer"] *,
    html body [data-testid="stChatMessageContent"],
    html body [data-testid="stChatMessageContent"] *,
    html body [data-testid="stSpinner"],
    html body [data-testid="stSpinner"] *,
    html body [data-testid="stSpinner"] div {{
        color: {WHITE} !important;
    }}
    html body .main .block-container a,
    html body .stMarkdown a {{
        color: {TEAL} !important;
    }}

    /* Keep button text in its intended color (overrides above) */
    html body .stButton > button,
    html body .stButton > button * {{
        color: {NAVY} !important;
    }}
    html body .stButton > button:hover,
    html body .stButton > button:hover * {{
        color: {NAVY} !important;
    }}

    /* Keep chat input text navy on its white background */
    html body [data-testid="stChatInput"] textarea,
    html body .stChatInput textarea {{
        color: {NAVY} !important;
    }}

    /* Custom classes */
    .main-header {{
        font-family: 'Poppins', sans-serif;
        font-size: 2rem;
        font-weight: 700;
        color: {WHITE};
        margin-bottom: 0.25rem;
    }}
    .sub-header {{
        font-family: 'Poppins', sans-serif;
        font-size: 1rem;
        color: #b8d4e3;
        margin-bottom: 1.5rem;
    }}
    .status-connected {{ color: {TEAL}; font-weight: 600; }}
    .status-error {{ color: #ff6b6b; font-weight: 600; }}

    /* Landing page styling */
    .landing-question {{
        font-family: 'Poppins', sans-serif;
        font-size: 2.2rem;
        font-weight: 600;
        color: {WHITE};
        text-align: center;
        margin: 2rem 0 2.5rem 0;
    }}
    .landing-subtitle {{
        font-family: 'Poppins', sans-serif;
        font-size: 1.05rem;
        font-weight: 300;
        color: #b8d4e3;
        text-align: center;
        margin-bottom: 3rem;
    }}

    /* Big landing buttons */
    .big-bot-button {{
        display: block;
        width: 100%;
        background-color: {TEAL};
        color: {NAVY} !important;
        font-family: 'Poppins', sans-serif;
        font-size: 1.4rem;
        font-weight: 600;
        text-align: center;
        padding: 2.5rem 1rem;
        border-radius: 14px;
        text-decoration: none !important;
        transition: all 0.25s ease;
        border: 2px solid {TEAL};
    }}
    .big-bot-button:hover {{
        background-color: {NAVY};
        color: {TEAL} !important;
        border: 2px solid {TEAL};
        transform: translateY(-3px);
    }}
    .big-bot-emoji {{
        font-size: 2.5rem;
        display: block;
        margin-bottom: 0.5rem;
    }}
    .big-bot-caption {{
        font-size: 0.9rem;
        font-weight: 400;
        margin-top: 0.5rem;
        opacity: 0.85;
    }}

    /* Hide the default Streamlit header/footer branding */
    #MainMenu {{ visibility: hidden; }}
    footer {{ visibility: hidden; }}
    header {{ visibility: hidden; }}
</style>
""", unsafe_allow_html=True)


# --- Logo Helper ---
def get_logo_base64():
    """Load the Ringover logo as base64 for inline embedding."""
    logo_path = os.path.join(os.path.dirname(__file__), "ringover-logo.png")
    if os.path.exists(logo_path):
        with open(logo_path, "rb") as f:
            return base64.b64encode(f.read()).decode("utf-8")
    return None


def show_logo(max_width_px=320):
    """Display the Ringover logo centered at the top."""
    logo_b64 = get_logo_base64()
    if logo_b64:
        st.markdown(
            f"""
            <div style="text-align: center; margin-top: 1.5rem; margin-bottom: 1rem;">
                <img src="data:image/png;base64,{logo_b64}" style="max-width: {max_width_px}px; width: 80%;" />
            </div>
            """,
            unsafe_allow_html=True,
        )
    else:
        # Fallback text if logo file not found
        st.markdown(
            f"""
            <div style="text-align: center; margin-top: 1.5rem; margin-bottom: 1rem;">
                <span style="color: {TEAL}; font-family: Poppins; font-weight: 700; font-size: 2.5rem;">ringover</span>
            </div>
            """,
            unsafe_allow_html=True,
        )


# --- Load Credentials ---
def load_credentials():
    """Load API keys and settings from Streamlit secrets or secrets.toml."""
    try:
        config = {
            "anthropic_api_key": st.secrets["secrets"]["ANTHROPIC_API_KEY"],
            # US Book of Business
            "google_sheet_id": st.secrets["secrets"].get("GOOGLE_SHEET_ID", ""),
            "google_sheet_gid": str(st.secrets["secrets"].get("GOOGLE_SHEET_GID", "")).strip(),
            "google_sheet_tab_name": st.secrets["secrets"].get("GOOGLE_SHEET_TAB_NAME", "").strip(),
            # UK Book of Business
            "uk_google_sheet_id": st.secrets["secrets"].get("UK_GOOGLE_SHEET_ID", ""),
            "uk_google_sheet_gid": str(st.secrets["secrets"].get("UK_GOOGLE_SHEET_GID", "")).strip(),
            "uk_google_sheet_tab_name": st.secrets["secrets"].get("UK_GOOGLE_SHEET_TAB_NAME", "").strip(),
            # Shared
            "google_drive_folder_id": st.secrets["secrets"].get("GOOGLE_DRIVE_FOLDER_ID", ""),
            "google_credentials": json.loads(st.secrets["secrets"]["GOOGLE_CREDENTIALS_JSON"]),
            "app_password": st.secrets["secrets"].get("APP_PASSWORD", ""),
        }
        return config
    except Exception as e:
        st.error(f"Could not load credentials. Make sure your secrets are set up correctly. Error: {e}")
        st.stop()


# --- Password Protection ---
def check_password(expected_password):
    """Returns True if user has entered the correct password."""
    if not expected_password:
        return True

    if st.session_state.get("authenticated", False):
        return True

    show_logo(max_width_px=260)
    st.markdown('<div class="landing-question" style="font-size:1.6rem;">Please enter the access password</div>', unsafe_allow_html=True)

    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        password = st.text_input("Password", type="password", key="password_input", label_visibility="collapsed", placeholder="Password")

        if password:
            if password == expected_password:
                st.session_state["authenticated"] = True
                st.rerun()
            else:
                st.error("❌ Incorrect password. Please try again.")

    return False


# --- Google Auth Helper ---
def get_google_creds(credentials_info):
    scopes = [
        "https://www.googleapis.com/auth/spreadsheets.readonly",
        "https://www.googleapis.com/auth/drive.readonly",
    ]
    return Credentials.from_service_account_info(credentials_info, scopes=scopes)


# ============================================================
# BOOK OF BUSINESS BOT
# ============================================================

@st.cache_data(ttl=300)
def load_sheet_data(_credentials, sheet_id, sheet_gid="", sheet_tab_name=""):
    """
    Loads rows from a specific tab of the Google Sheet.

    Priority:
      1. sheet_gid    (numeric tab id from the URL, e.g. gid=1485115496)
      2. sheet_tab_name (human-readable tab name)
      3. first tab (fallback, for backward compatibility)
    """
    creds = get_google_creds(_credentials)
    client = gspread.authorize(creds)
    spreadsheet = client.open_by_key(sheet_id)

    worksheet = None
    if sheet_gid:
        try:
            gid_int = int(sheet_gid)
            for ws in spreadsheet.worksheets():
                if ws.id == gid_int:
                    worksheet = ws
                    break
        except (ValueError, TypeError):
            pass

    if worksheet is None and sheet_tab_name:
        try:
            worksheet = spreadsheet.worksheet(sheet_tab_name)
        except Exception:
            worksheet = None

    if worksheet is None:
        worksheet = spreadsheet.sheet1

    data = worksheet.get_all_records()
    return data, worksheet.title, [w.title for w in spreadsheet.worksheets()]


def data_to_dataframe(data):
    """Convert the sheet rows into a pandas DataFrame with numeric columns normalized."""
    if not data:
        return pd.DataFrame()

    df = pd.DataFrame(data)

    # Try to coerce columns that look like currency/numbers to float.
    # A column is numeric-looking if > 70% of non-blank values match a currency/number pattern.
    number_pattern = re.compile(r"^-?\$?\s*-?[\d,]+(\.\d+)?%?$")
    for col in df.columns:
        if df[col].dtype != object:
            continue
        non_blank = df[col].astype(str).str.strip()
        non_blank = non_blank[non_blank != ""]
        if len(non_blank) == 0:
            continue
        matches = non_blank.apply(lambda v: bool(number_pattern.match(str(v).strip())))
        if matches.mean() > 0.7:
            cleaned = (
                df[col]
                .astype(str)
                .str.replace("$", "", regex=False)
                .str.replace(",", "", regex=False)
                .str.replace("%", "", regex=False)
                .str.strip()
                .replace("", None)
            )
            try:
                df[col] = pd.to_numeric(cleaned, errors="coerce")
            except Exception:
                pass

    return df


def dataframe_to_schema(df):
    """Return a string describing the DataFrame columns, types, and sample values for Claude."""
    if df.empty:
        return "The DataFrame is empty."

    lines = [
        f"DataFrame name: df",
        f"Total rows: {len(df)}",
        f"Total columns: {len(df.columns)}",
        "",
        "Columns (name | dtype | non-null count | sample values):",
    ]
    for col in df.columns:
        dtype = str(df[col].dtype)
        non_null = int(df[col].notna().sum())
        sample = df[col].dropna().astype(str).head(4).tolist()
        sample_str = ", ".join(f'"{s}"' for s in sample)
        lines.append(f'  - "{col}" | {dtype} | {non_null} non-null | samples: [{sample_str}]')

    return "\n".join(lines)


BOB_SYSTEM_PROMPT = """You are the Ringover Book of Business Bot — a helpful assistant for the Ringover US team.

You have access to a pandas DataFrame named `df` that holds the complete US Book of Business data
(roughly {row_count} rows). You must use the `query_dataframe` tool to compute ANY answer that involves:
- counts, sums, averages, or other aggregations
- filters across the data (e.g., "accounts where X and Y")
- sorting or ranking (e.g., "top N by MRR")
- looking up specific accounts or values

NEVER try to answer such questions from memory or by reading sample values — always run a query first.
You may run MULTIPLE queries if you need intermediate results.

COLUMN NAMES ARE CASE-SENSITIVE and must be quoted exactly as they appear in the schema below.
If a user's wording doesn't exactly match a column name, pick the column whose name is the closest
semantic match (for example, "Account Director" might be a value inside a column like "Upsell Owner"
or "Owner Type" — scan the samples to find the right fit, and run a small exploratory query if unsure).

When comparing text values, be mindful of whitespace and case. You may use `.str.strip()` and
`.str.lower()` for fuzzier matching when appropriate.

When presenting results to the user:
- Lead with the direct answer (e.g., "There are 324 accounts matching that criteria.")
- When listing accounts, show 10–20 by default and offer to expand if the list is longer
- Include relevant context columns (account name, MRR, CSM, etc.)
- Format money, counts, and percentages clearly
- If a query returns zero rows, say so plainly and suggest how the question might be rephrased
- Always cite the exact numbers from query results — never round or estimate

DataFrame schema:
{schema}
"""


QUERY_TOOL = {
    "name": "query_dataframe",
    "description": (
        "Evaluate a Python expression against the Book of Business pandas DataFrame. "
        "The DataFrame is available as `df`. `pd` (pandas) is also available. "
        "Use this for ALL counts, filters, sums, sorts, and lookups. "
        "Return value is converted to a readable string or JSON and sent back to you. "
        "Examples of valid expressions:\n"
        "  - len(df)\n"
        "  - df['Total MRR'].sum()\n"
        "  - len(df[(df['Upsell Owner'] == 'Account Director') & (df['Empower MRR'] == 0)])\n"
        "  - df.nlargest(20, 'Total MRR')[['Account Name', 'Total MRR', 'CSM']].to_dict('records')\n"
        "  - df['CSM'].value_counts().to_dict()\n"
        "  - df[df['Account Name'].str.contains('Acme', case=False, na=False)].to_dict('records')"
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "expression": {
                "type": "string",
                "description": (
                    "A single Python expression to evaluate. The DataFrame is `df` and pandas is `pd`. "
                    "Do NOT use statements (no assignments, no imports, no print). Just an expression."
                ),
            }
        },
        "required": ["expression"],
    },
}


def run_pandas_query(df, expression):
    """Safely evaluate a pandas expression and return a string result."""
    try:
        # Restricted eval — only df, pd, and a small set of safe builtins
        safe_builtins = {
            "len": len, "sum": sum, "min": min, "max": max, "sorted": sorted,
            "round": round, "abs": abs, "int": int, "float": float, "str": str,
            "list": list, "dict": dict, "set": set, "tuple": tuple, "bool": bool,
            "range": range, "enumerate": enumerate, "zip": zip, "any": any, "all": all,
        }
        result = eval(expression, {"__builtins__": safe_builtins}, {"df": df, "pd": pd})
    except Exception as e:
        return f"ERROR running expression `{expression}`: {type(e).__name__}: {e}"

    # Serialize the result so Claude can read it
    try:
        if isinstance(result, pd.DataFrame):
            # Limit row count in output to avoid blowing context
            if len(result) > 50:
                head = result.head(50).to_dict("records")
                return json.dumps({"total_rows": len(result), "showing_first_50": head}, default=str)
            return json.dumps(result.to_dict("records"), default=str)
        elif isinstance(result, pd.Series):
            if len(result) > 100:
                head = result.head(100).to_dict()
                return json.dumps({"total_items": len(result), "showing_first_100": head}, default=str)
            return json.dumps(result.to_dict(), default=str)
        elif isinstance(result, (int, float, str, bool)) or result is None:
            return json.dumps(result, default=str)
        elif isinstance(result, (list, dict, tuple, set)):
            if isinstance(result, set):
                result = list(result)
            return json.dumps(result, default=str)
        else:
            return str(result)
    except Exception as e:
        return f"Result computed but could not be serialized: {type(e).__name__}: {e}. repr: {repr(result)[:500]}"


# ============================================================
# PROCESS BOT
# ============================================================

# --- Helpers to read non-Google-native file types ---

def _download_file_bytes(service, file_id):
    """Download a non-Google-native file from Drive and return raw bytes."""
    request = service.files().get_media(fileId=file_id)
    buffer = io.BytesIO()
    downloader = MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    buffer.seek(0)
    return buffer


def _read_pdf(buffer):
    """Extract text from a PDF file buffer."""
    try:
        import pypdf
        reader = pypdf.PdfReader(buffer)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages) if pages else "[PDF contained no extractable text]"
    except Exception as e:
        return f"[Could not read PDF: {e}]"


def _read_docx(buffer):
    """Extract text from a Word .docx file buffer."""
    try:
        import docx
        doc = docx.Document(buffer)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also read tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                paragraphs.append(" | ".join(cells))
        return "\n".join(paragraphs) if paragraphs else "[Word document was empty]"
    except Exception as e:
        return f"[Could not read Word document: {e}]"


def _read_pptx(buffer):
    """Extract text from a PowerPoint .pptx file buffer."""
    try:
        from pptx import Presentation
        prs = Presentation(buffer)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if texts:
                slides_text.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides_text) if slides_text else "[Presentation was empty]"
    except Exception as e:
        return f"[Could not read PowerPoint: {e}]"


def _read_xlsx(buffer):
    """Extract text from an Excel .xlsx file buffer."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(buffer, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        return "\n\n".join(sheets_text) if sheets_text else "[Spreadsheet was empty]"
    except Exception as e:
        return f"[Could not read Excel file: {e}]"


def _read_uploaded_file(service, file_id, name, mime):
    """Read a non-Google-native file uploaded to Drive (PDF, docx, pptx, xlsx, etc.)."""
    try:
        buffer = _download_file_bytes(service, file_id)
    except Exception as e:
        return f"[Could not download file: {e}]"

    lower_name = name.lower()

    # PDF
    if mime == "application/pdf" or lower_name.endswith(".pdf"):
        return _read_pdf(buffer)

    # Word
    if mime in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword") or lower_name.endswith((".docx", ".doc")):
        return _read_docx(buffer)

    # PowerPoint
    if mime in ("application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint") or lower_name.endswith((".pptx", ".ppt")):
        return _read_pptx(buffer)

    # Excel
    if mime in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel") or lower_name.endswith((".xlsx", ".xls")):
        return _read_xlsx(buffer)

    # Plain text fallback
    try:
        buffer.seek(0)
        return buffer.read().decode("utf-8", errors="replace")
    except Exception:
        return f"[Unsupported file type: {mime}]"


# --- Load all documents from Drive (simple approach — all docs in every prompt) ---

@st.cache_data(ttl=600)
def load_drive_documents(_credentials, folder_id):
    """Load ALL documents from Drive with full content.
    Supports Google-native formats + PDF, Word, PowerPoint, Excel."""
    creds = get_google_creds(_credentials)
    service = build("drive", "v3", credentials=creds)

    documents = []

    def read_folder(fid):
        page_token = None
        while True:
            query = f"'{fid}' in parents and trashed = false"
            results = service.files().list(
                q=query,
                fields="nextPageToken, files(id, name, mimeType)",
                pageSize=100,
                pageToken=page_token,
            ).execute()
            files = results.get("files", [])

            for file in files:
                mime = file["mimeType"]
                name = file["name"]
                file_id = file["id"]

                if mime == "application/vnd.google-apps.folder":
                    read_folder(file_id)
                    continue

                content = None

                # Google Docs
                if mime == "application/vnd.google-apps.document":
                    try:
                        raw = service.files().export(fileId=file_id, mimeType="text/plain").execute()
                        content = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                    except Exception:
                        content = "[Could not read this document]"

                # Google Sheets
                elif mime == "application/vnd.google-apps.spreadsheet":
                    try:
                        raw = service.files().export(fileId=file_id, mimeType="text/csv").execute()
                        content = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                    except Exception:
                        content = "[Could not read this spreadsheet]"

                # Google Slides
                elif mime == "application/vnd.google-apps.presentation":
                    try:
                        raw = service.files().export(fileId=file_id, mimeType="text/plain").execute()
                        content = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                    except Exception:
                        content = "[Could not read this presentation]"

                # Plain text / CSV / Markdown
                elif mime in ("text/plain", "text/csv", "text/markdown"):
                    try:
                        buf = _download_file_bytes(service, file_id)
                        content = buf.read().decode("utf-8", errors="replace")
                    except Exception:
                        content = "[Could not read this file]"

                # PDF, Word, PowerPoint, Excel (uploaded files)
                elif (mime.startswith("application/pdf")
                      or mime.startswith("application/vnd.openxmlformats")
                      or mime.startswith("application/msword")
                      or mime.startswith("application/vnd.ms-")
                      or name.lower().endswith((".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"))):
                    content = _read_uploaded_file(service, file_id, name, mime)

                # Shortcuts
                elif mime == "application/vnd.google-apps.shortcut":
                    try:
                        details = service.files().get(fileId=file_id, fields="shortcutDetails").execute()
                        target_id = details.get("shortcutDetails", {}).get("targetId")
                        target_mime = details.get("shortcutDetails", {}).get("targetMimeType", "")
                        if target_id:
                            if target_mime == "application/vnd.google-apps.document":
                                raw = service.files().export(fileId=target_id, mimeType="text/plain").execute()
                                content = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                            elif target_mime == "application/vnd.google-apps.folder":
                                read_folder(target_id)
                                continue
                            elif target_mime == "application/vnd.google-apps.spreadsheet":
                                raw = service.files().export(fileId=target_id, mimeType="text/csv").execute()
                                content = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                            elif target_mime == "application/vnd.google-apps.presentation":
                                raw = service.files().export(fileId=target_id, mimeType="text/plain").execute()
                                content = raw.decode("utf-8") if isinstance(raw, bytes) else raw
                            else:
                                content = _read_uploaded_file(service, target_id, name, target_mime)
                    except Exception:
                        content = "[Could not resolve this shortcut]"

                if content is not None:
                    documents.append({"name": name, "content": content})

            page_token = results.get("nextPageToken")
            if not page_token:
                break

    read_folder(folder_id)
    return documents


def _score_document(doc, keywords):
    """Score a document's relevance to a question based on keyword matches.
    Returns a numeric score — higher is more relevant."""
    name_lower = doc["name"].lower()
    content_lower = doc["content"].lower()
    score = 0

    for kw in keywords:
        # Matches in document name are worth more (title relevance)
        if kw in name_lower:
            score += 10
        # Count occurrences in content (capped to avoid one keyword dominating)
        count = content_lower.count(kw)
        score += min(count, 20)

    return score


def _extract_keywords(question):
    """Extract meaningful keywords from a question for document matching."""
    # Common words to ignore
    stop_words = {
        "a", "an", "the", "is", "are", "was", "were", "be", "been", "being",
        "have", "has", "had", "do", "does", "did", "will", "would", "could",
        "should", "may", "might", "shall", "can", "need", "dare", "ought",
        "used", "to", "of", "in", "for", "on", "with", "at", "by", "from",
        "as", "into", "through", "during", "before", "after", "above", "below",
        "between", "out", "off", "over", "under", "again", "further", "then",
        "once", "here", "there", "when", "where", "why", "how", "all", "both",
        "each", "few", "more", "most", "other", "some", "such", "no", "nor",
        "not", "only", "own", "same", "so", "than", "too", "very", "just",
        "about", "what", "which", "who", "whom", "this", "that", "these",
        "those", "am", "if", "or", "and", "but", "because", "until", "while",
        "it", "its", "i", "me", "my", "we", "our", "you", "your", "he", "she",
        "they", "them", "his", "her", "us", "tell", "show", "give", "get",
        "find", "know", "think", "see", "want", "look", "make", "go", "say",
        "let", "please", "also", "like", "much", "many", "well", "still",
    }

    # Split on non-alphanumeric, lowercase, remove short/stop words
    words = re.split(r"[^a-zA-Z0-9]+", question.lower())
    keywords = [w for w in words if len(w) > 2 and w not in stop_words]

    # Also include 2-word and 3-word phrases from the original question for better matching
    q_lower = question.lower()
    words_raw = re.split(r"[^a-zA-Z0-9]+", q_lower)
    words_raw = [w for w in words_raw if w]
    for i in range(len(words_raw) - 1):
        phrase = words_raw[i] + " " + words_raw[i + 1]
        if words_raw[i] not in stop_words or words_raw[i + 1] not in stop_words:
            keywords.append(phrase)
    for i in range(len(words_raw) - 2):
        phrase = words_raw[i] + " " + words_raw[i + 1] + " " + words_raw[i + 2]
        keywords.append(phrase)

    return keywords


def filter_relevant_documents(documents, question, chat_history=None, max_docs=15):
    """Filter documents by keyword relevance to the question.
    Also considers recent chat history so follow-ups pull the same docs."""

    # Build keyword list from the current question
    keywords = _extract_keywords(question)

    # Also extract keywords from recent conversation for follow-up context
    if chat_history:
        recent = chat_history[-4:]  # last 2 Q&A pairs
        for msg in recent:
            if msg["role"] == "user":
                keywords.extend(_extract_keywords(msg["content"]))
            elif msg["role"] == "assistant":
                # Extract any document names mentioned in previous answers
                for doc in documents:
                    if doc["name"].lower() in msg["content"].lower():
                        # Boost this doc heavily by adding its name words as keywords
                        name_words = re.split(r"[^a-zA-Z0-9]+", doc["name"].lower())
                        keywords.extend([w for w in name_words if len(w) > 2])

    if not keywords:
        # No meaningful keywords — return first max_docs documents
        return documents[:max_docs]

    # Score and rank all documents
    scored = [(doc, _score_document(doc, keywords)) for doc in documents]
    scored.sort(key=lambda x: x[1], reverse=True)

    # Take top matches (only those with score > 0)
    relevant = [doc for doc, score in scored[:max_docs] if score > 0]

    if not relevant:
        # No keyword matches at all — return first max_docs as fallback
        return documents[:max_docs]

    return relevant


def documents_to_summary(documents):
    """Build a text summary of the provided documents for the system prompt."""
    if not documents:
        return "No documents found in the Google Drive folder."

    MAX_PER_DOC = 80000
    MAX_TOTAL = 200000

    summary = f"Providing {len(documents)} documents:\n\n"
    total_chars = len(summary)

    for i, doc in enumerate(documents):
        header = f"--- DOCUMENT: {doc['name']} ---\n"
        content = doc["content"]

        if len(content) > MAX_PER_DOC:
            content = content[:MAX_PER_DOC] + "\n... [Document truncated due to length]"

        section = header + content + "\n\n"

        if total_chars + len(section) > MAX_TOTAL:
            remaining = MAX_TOTAL - total_chars
            if remaining > len(header) + 200:
                section = header + content[:remaining - len(header) - 60] + "\n... [Truncated to fit context limit]\n\n"
                summary += section
            break

        summary += section
        total_chars += len(section)

    return summary


PROCESS_SYSTEM_PROMPT = """You are the Ringover Process Bot — a helpful assistant for the Ringover US team.
You have access to process documents from the team's Google Drive folder.

Your job is to answer questions about company processes, procedures, and guidelines by finding
the answer in the documents provided below.

CRITICAL RULES — YOU MUST FOLLOW THESE:
1. ONLY answer using information found in the documents below. If the answer is not in these documents, say "I couldn't find that information in the available documents" — do NOT guess or make up an answer.
2. When you find an answer, ALWAYS cite the exact document name (e.g., "According to 'Tips and Tricks for creating a purchase order'...").
3. Quote or closely paraphrase the relevant text from the document so the user can verify your answer.
4. If the answer spans multiple documents, reference all of them.
5. If a process has specific steps, list them clearly in order exactly as they appear in the document.
6. If you notice conflicting information between documents, flag it to the user.
7. Be conversational and helpful. These are your colleagues.
8. NEVER fabricate process steps, terms, definitions, or procedures that are not explicitly written in the documents.

HERE ARE THE CURRENT PROCESS DOCUMENTS:
{data}
"""


# ============================================================
# SHARED: Ask Claude
# ============================================================

def ask_claude(client, question, system_prompt, chat_history):
    """Simple Claude call without tools — used by Process Bot."""
    messages = []
    for msg in chat_history:
        messages.append({"role": msg["role"], "content": msg["content"]})
    messages.append({"role": "user", "content": question})

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4096,
        system=system_prompt,
        messages=messages,
    )

    return response.content[0].text


def ask_claude_with_dataframe(client, question, system_prompt, chat_history, df, max_iterations=8):
    """Claude call with access to query_dataframe tool — used by Book of Business Bot."""
    # Rebuild message history as plain text messages (tool traffic is not persisted across turns)
    messages = []
    for msg in chat_history:
        messages.append({"role": msg["role"], "content": msg["content"]})
    messages.append({"role": "user", "content": question})

    for _ in range(max_iterations):
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4096,
            system=system_prompt,
            tools=[QUERY_TOOL],
            messages=messages,
        )

        if response.stop_reason == "tool_use":
            # Append assistant's tool-call message
            messages.append({"role": "assistant", "content": response.content})

            # Run every tool call in this turn and append results
            tool_results = []
            for block in response.content:
                if getattr(block, "type", None) == "tool_use" and block.name == "query_dataframe":
                    expression = block.input.get("expression", "")
                    result_text = run_pandas_query(df, expression)
                    # Truncate very long results to keep context manageable
                    if len(result_text) > 20000:
                        result_text = result_text[:20000] + "\n... [truncated]"
                    tool_results.append({
                        "type": "tool_result",
                        "tool_use_id": block.id,
                        "content": result_text,
                    })
            messages.append({"role": "user", "content": tool_results})
            continue

        # Final response — pull out any text blocks
        final_text = ""
        for block in response.content:
            if getattr(block, "type", None) == "text":
                final_text += block.text
        return final_text if final_text.strip() else "I couldn't produce an answer. Try rephrasing the question."

    return "I ran too many queries without reaching a final answer. Try simplifying the question or breaking it into parts."


# ============================================================
# LANDING PAGE
# ============================================================

def show_landing_page():
    show_logo(max_width_px=340)

    st.markdown('<div class="landing-question">Which bot do you need to talk to?</div>', unsafe_allow_html=True)
    st.markdown('<div class="landing-subtitle">Choose a bot below to get started.</div>', unsafe_allow_html=True)

    col_left, col_us, col_gap1, col_uk, col_gap2, col_process, col_right = st.columns([0.5, 3, 0.3, 3, 0.3, 3, 0.5])

    # Load flag images as base64 (same approach as the Ringover logo)
    flag_dir = os.path.dirname(__file__)
    def _load_flag_b64(filename):
        fpath = os.path.join(flag_dir, filename)
        if os.path.exists(fpath):
            with open(fpath, "rb") as f:
                return base64.b64encode(f.read()).decode("utf-8")
        return None

    us_flag_b64 = _load_flag_b64("flag_us.png")
    gb_flag_b64 = _load_flag_b64("flag_gb.png")

    def _flag_html(b64):
        if b64:
            return f'<div style="text-align:center; margin-bottom:0.5rem;"><img src="data:image/png;base64,{b64}" style="width:56px; border-radius:4px;" /></div>'
        return ""

    with col_us:
        st.markdown(_flag_html(us_flag_b64), unsafe_allow_html=True)
        if st.button("USA BoB Bot\n\nUS Book of Business", key="goto_bob", use_container_width=True):
            st.session_state["view"] = "bob"
            st.rerun()
        st.markdown(
            '<div style="text-align:center; color:#b8d4e3; font-size:0.9rem; margin-top:0.75rem;">'
            'US accounts, MRR, integrations, renewals'
            '</div>',
            unsafe_allow_html=True,
        )

    with col_uk:
        st.markdown(_flag_html(gb_flag_b64), unsafe_allow_html=True)
        if st.button("UK BoB Bot\n\nUK Book of Business", key="goto_uk_bob", use_container_width=True):
            st.session_state["view"] = "uk_bob"
            st.rerun()
        st.markdown(
            '<div style="text-align:center; color:#b8d4e3; font-size:0.9rem; margin-top:0.75rem;">'
            'UK accounts, MRR, integrations, renewals'
            '</div>',
            unsafe_allow_html=True,
        )

    with col_process:
        st.markdown(_flag_html(us_flag_b64), unsafe_allow_html=True)
        if st.button("Process Bot\n\nCompany Processes", key="goto_process", use_container_width=True):
            st.session_state["view"] = "process"
            st.rerun()
        st.markdown(
            '<div style="text-align:center; color:#b8d4e3; font-size:0.9rem; margin-top:0.75rem;">'
            'Procedures, guidelines, how-tos'
            '</div>',
            unsafe_allow_html=True,
        )


# ============================================================
# BOT VIEW (shared logic for BoB + Process)
# ============================================================

def show_bot_view(bot_mode, config):
    # --- Load data BEFORE rendering sidebar so we know what to display ---
    data = None
    documents = None
    df = None
    connection_ok = False
    connection_error = None

    if bot_mode in ("bob", "uk_bob"):
        if bot_mode == "uk_bob":
            sheet_id = config.get("uk_google_sheet_id", "")
            sheet_gid = config.get("uk_google_sheet_gid", "")
            sheet_tab = config.get("uk_google_sheet_tab_name", "")
            currency_symbol = "£"
            region_label = "UK"
        else:
            sheet_id = config.get("google_sheet_id", "")
            sheet_gid = config.get("google_sheet_gid", "")
            sheet_tab = config.get("google_sheet_tab_name", "")
            currency_symbol = "$"
            region_label = "US"

        try:
            data, loaded_tab, all_tabs = load_sheet_data(
                config["google_credentials"],
                sheet_id,
                sheet_gid,
                sheet_tab,
            )
            df = data_to_dataframe(data)
            connection_ok = True
        except Exception as e:
            connection_error = str(e)
    else:
        region_label = ""
        currency_symbol = ""
        loaded_tab = ""
        all_tabs = []
        try:
            documents = load_drive_documents(config["google_credentials"], config["google_drive_folder_id"])
            connection_ok = True
        except Exception as e:
            connection_error = str(e)

    # --- Sidebar (always renders regardless of connection status) ---
    with st.sidebar:
        show_logo(max_width_px=180)

        if st.button("← Back to Home", key="back_home", use_container_width=True):
            st.session_state["view"] = "landing"
            st.rerun()

        st.markdown("---")
        st.markdown("### 📡 Connection Status")

        if bot_mode in ("bob", "uk_bob"):
            if connection_ok and df is not None:
                st.markdown(f'<span class="status-connected">✅ Connected to {region_label} Google Sheet</span>', unsafe_allow_html=True)
                st.markdown(f"**Active tab:** {loaded_tab}")
                if len(all_tabs) > 1:
                    st.caption(f"All tabs in this sheet: {', '.join(all_tabs)}")
                st.markdown(f"**Accounts loaded:** {len(df)}")
                st.markdown(f"**Data columns:** {len(df.columns)}")

                mrr_col = next((c for c in df.columns if "total mrr" in c.lower()), None)
                if mrr_col and pd.api.types.is_numeric_dtype(df[mrr_col]):
                    total_mrr = df[mrr_col].sum()
                    st.markdown(f"**Total MRR:** {currency_symbol}{total_mrr:,.2f}")
            else:
                st.markdown(f'<span class="status-error">❌ Could not connect to {region_label} Google Sheet</span>', unsafe_allow_html=True)
                if connection_error:
                    st.error(f"Error: {connection_error}")
        else:
            if connection_ok and documents is not None:
                st.markdown(f'<span class="status-connected">✅ Connected to Google Drive</span>', unsafe_allow_html=True)
                st.markdown(f"**Documents loaded:** {len(documents)}")
                if documents:
                    st.markdown("**Documents found:**")
                    for doc in documents:
                        st.markdown(f"- {doc['name']}")
            else:
                st.markdown(f'<span class="status-error">❌ Could not connect to Google Drive</span>', unsafe_allow_html=True)
                if connection_error:
                    st.error(f"Error: {connection_error}")

        st.markdown("---")
        if st.button("🔄 Refresh Data", use_container_width=True):
            st.cache_data.clear()
            st.rerun()

        st.markdown("---")
        if bot_mode == "bob":
            st.markdown("### 💡 Example Questions")
            st.markdown("""
- Which accounts have Bullhorn integrated?
- List all platinum accounts with their MRR
- Who are the top 10 accounts by MRR?
- Which accounts are up for renewal?
- How many accounts does each CSM manage?
- What is the total MRR by business unit?
- Show me all accounts with Empower licenses
            """)
        elif bot_mode == "uk_bob":
            st.markdown("### 💡 Example Questions")
            st.markdown("""
- Who are the top 10 UK accounts by MRR?
- Which UK accounts have an ATS/CRM integrated?
- List all UK accounts up for renewal
- What is the total UK MRR?
- Which accounts have Empower licences?
- Show me all accounts with whitespace opportunities
- How many accounts does each team manage?
            """)
        else:
            st.markdown("### 💡 Example Questions")
            st.markdown("""
- What is the process for client onboarding?
- How do we handle an escalation?
- What are the steps for a renewal?
- What is our SLA for support tickets?
- How do I submit a feature request?
            """)

    # Main area header
    if bot_mode == "bob":
        st.markdown('<div class="main-header">USA Book of Business Bot</div>', unsafe_allow_html=True)
        st.markdown('<div class="sub-header">Ask me anything about the US Book of Business — accounts, MRR, integrations, renewals, and more.</div>', unsafe_allow_html=True)
    elif bot_mode == "uk_bob":
        st.markdown('<div class="main-header">UK Book of Business Bot</div>', unsafe_allow_html=True)
        st.markdown('<div class="sub-header">Ask me anything about the UK Book of Business — accounts, MRR, integrations, renewals, and more.</div>', unsafe_allow_html=True)
    else:
        st.markdown('<div class="main-header">📋 Process Bot</div>', unsafe_allow_html=True)
        st.markdown('<div class="sub-header">Ask me anything about Ringover processes, procedures, and guidelines.</div>', unsafe_allow_html=True)

    # Chat history
    bob_key = "bob_messages"
    uk_bob_key = "uk_bob_messages"
    process_key = "process_messages"

    if bob_key not in st.session_state:
        st.session_state[bob_key] = []
    if uk_bob_key not in st.session_state:
        st.session_state[uk_bob_key] = []
    if process_key not in st.session_state:
        st.session_state[process_key] = []

    if bot_mode == "bob":
        current_key = bob_key
    elif bot_mode == "uk_bob":
        current_key = uk_bob_key
    else:
        current_key = process_key

    for message in st.session_state[current_key]:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    if bot_mode == "bob":
        placeholder = "Ask me about the US Book of Business..."
    elif bot_mode == "uk_bob":
        placeholder = "Ask me about the UK Book of Business..."
    else:
        placeholder = "Ask me about a process or procedure..."

    # Show a warning in the main area if the data connection failed
    if not connection_ok:
        st.warning(f"⚠️ Could not connect to the data source. Please check your secrets configuration and click **Refresh Data** in the sidebar to retry.")

    if prompt := st.chat_input(placeholder):
        if not connection_ok:
            st.error("Cannot process questions without a data connection. Please fix the connection and refresh.")
            st.stop()
        with st.chat_message("user"):
            st.markdown(prompt)
        st.session_state[current_key].append({"role": "user", "content": prompt})

        with st.chat_message("assistant"):
            with st.spinner("Querying the data..." if bot_mode in ("bob", "uk_bob") else "Searching the documents..."):
                try:
                    claude_client = anthropic.Anthropic(api_key=config["anthropic_api_key"])

                    if bot_mode in ("bob", "uk_bob"):
                        schema = dataframe_to_schema(df)
                        if bot_mode == "uk_bob":
                            system = BOB_SYSTEM_PROMPT.replace("US team", "UK team").replace("US Book of Business", "UK Book of Business").format(row_count=len(df), schema=schema)
                        else:
                            system = BOB_SYSTEM_PROMPT.format(row_count=len(df), schema=schema)
                        response = ask_claude_with_dataframe(
                            claude_client,
                            prompt,
                            system,
                            st.session_state[current_key][:-1],
                            df,
                        )
                    else:
                        relevant_docs = filter_relevant_documents(
                            documents,
                            prompt,
                            chat_history=st.session_state[current_key][:-1],
                        )
                        doc_summary = documents_to_summary(relevant_docs)
                        system = PROCESS_SYSTEM_PROMPT.format(data=doc_summary)
                        response = ask_claude(
                            claude_client,
                            prompt,
                            system,
                            st.session_state[current_key][:-1],
                        )

                    st.markdown(response)
                    st.session_state[current_key].append({"role": "assistant", "content": response})
                except anthropic.AuthenticationError:
                    st.error("Your Anthropic API key is invalid. Please check your secrets.")
                except Exception as e:
                    st.error(f"Something went wrong: {e}")


# ============================================================
# MAIN APP
# ============================================================

def main():
    config = load_credentials()

    if not check_password(config["app_password"]):
        st.stop()

    # Track which view we're on
    if "view" not in st.session_state:
        st.session_state["view"] = "landing"

    view = st.session_state["view"]

    if view == "landing":
        show_landing_page()
    elif view == "bob":
        show_bot_view("bob", config)
    elif view == "uk_bob":
        show_bot_view("uk_bob", config)
    elif view == "process":
        show_bot_view("process", config)
    else:
        st.session_state["view"] = "landing"
        st.rerun()


if __name__ == "__main__":
    main()
